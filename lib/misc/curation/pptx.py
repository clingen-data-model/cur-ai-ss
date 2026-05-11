import io
from pathlib import Path
from typing import TYPE_CHECKING

from pptx import Presentation
from pptx.util import Inches, Pt

from lib.misc.curation.models import CurationSummaryRow, SectionContent

if TYPE_CHECKING:
    from pptx.text.text import TextFrame

# python-pptx API Guide:
# - Inches(n): Converts inches to EMUs (English Metric Units). PowerPoint uses EMUs internally.
#   1 inch = 914,400 EMUs. Inches() is a convenience function for human-readable measurements.
# - Pt(n): Converts points to EMUs for font sizes. 1 point ≈ 12,700 EMUs.
# - Presentation(): Creates a new presentation object (in-memory, 16:9 aspect by default)
# - prs.slide_width/slide_height: Set slide dimensions in EMUs
# - prs.slide_layouts[i]: Accesses predefined layouts (11 total, 0-10):
#   0=Title Slide, 1=Title and Content, 2=Section Header, 3=Two Content, 4=Comparison,
#   5=Title Only, 6=Blank, 7=Content with Caption, 8=Picture with Caption,
#   9=Title and Vertical Text, 10=Vertical Title and Text
# - slide.shapes: Container for all objects on a slide (tables, text boxes, images, etc.)
# - TextFrame: Represents text within a shape; allows adding paragraphs and formatting


def _format_cell_with_sections(
    text_frame: 'TextFrame', sections: list[SectionContent]
) -> None:
    """Format a cell with bold section titles and paragraph breaks.

    Uses the TextFrame API to add paragraphs and runs with different formatting.
    Demonstrates:
    - text_frame.clear(): Removes all existing paragraphs
    - text_frame.add_paragraph(): Adds a new paragraph (equivalent to pressing Enter)
    - paragraph.add_run(): Adds a run (contiguous text with same formatting)
    - run.font properties: Modify character-level formatting (bold, size, color, etc.)
    """
    # Clear existing content from the text frame. This ensures we start fresh.
    text_frame.clear()

    for idx, section in enumerate(sections):
        # Add a new paragraph to the text frame. Paragraphs represent logical text blocks.
        # Each paragraph has independent line spacing, alignment, indentation properties.
        p = text_frame.add_paragraph()

        # Add a run for the section title. A run is the smallest unit of consistent formatting.
        # Multiple runs within the same paragraph can have different formatting.
        title_run = p.add_run()
        title_run.text = section.title
        # Make the title bold using the font property of the run.
        title_run.font.bold = True

        # Add the content text if present. Prefix with newline for visual separation.
        if section.content:
            # Add a second run to the same paragraph for the content (different formatting).
            content_run = p.add_run()
            content_run.text = '\n' + section.content

        # Add a blank paragraph between sections for visual spacing (except after the last section).
        if idx < len(sections) - 1:
            text_frame.add_paragraph()


def build_curation_pptx(rows: list[CurationSummaryRow]) -> bytes:
    """Generate a PPTX slide with curation summary table.

    Creates a single landscape slide with a table containing:
    - Header row: Publication & Testing | Proband | Variant Info | Clinical Presentation | Functional/Segregation | Score
    - One data row per CurationSummaryRow
    - Optional pedigree image if available

    Returns raw PPTX bytes.
    """
    # Create an in-memory Presentation object. Each presentation has a default 16:9 aspect ratio.
    prs = Presentation()
    # Set slide dimensions in inches (converted to EMUs internally by python-pptx).
    # 13" x 7.5" gives us a landscape layout suitable for wide tables.
    prs.slide_width = Inches(13)
    prs.slide_height = Inches(7.5)

    # Retrieve the Blank layout (index 6) to avoid default title/content placeholders.
    # This gives us full control over shape placement.
    blank_layout = prs.slide_layouts[6]
    # Add a new slide using the blank layout. prs.slides is a collection that supports add_slide().
    slide = prs.slides.add_slide(blank_layout)

    # Table dimensions: (rows + 1 header) x 6 columns
    rows_count = len(rows) + 1  # +1 for header
    cols_count = 6

    # Define table position and size in EMUs (via Inches conversion).
    # Position: 0.5" from left and top edges. Size: 12" wide x 6.5" tall.
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(12)
    height = Inches(6.5)

    # Add a table shape to the slide. Returns a GraphicsFrame that contains the table.
    # slide.shapes.add_table(rows, cols, left, top, width, height) creates a new table
    # at the specified position and size. The shape object wraps a Table object.
    table_shape = slide.shapes.add_table(
        rows_count, cols_count, left, top, width, height
    )
    # Extract the Table object from the shape. Use this to access rows, columns, and cells.
    table = table_shape.table

    # Distribute column widths equally. Note: width is an EMU value from Inches(12).
    # Divide total width by column count to get uniform column widths.
    col_width = int(width / cols_count)
    for col_idx in range(cols_count):
        # Access columns by index and set their width property (in EMUs).
        table.columns[col_idx].width = col_width

    # Header row - set smaller height to save vertical space.
    # table.rows[i] accesses rows by index (0 is header, 1+ are data rows).
    header_row = table.rows[0]
    # Set row height in EMUs via Inches conversion. 0.4" is compact for header.
    header_row.height = Inches(0.4)

    headers = [
        'Publication & Testing',
        'Proband',
        'Variant Info',
        'Clinical Presentation',
        'Functional/Segregation',
        'Score',
    ]
    for col_idx, header_text in enumerate(headers):
        # Access cell by row and column indices. table.cell(row, col) returns a Cell object.
        cell = table.cell(0, col_idx)
        # Set plain text directly. Note: cell.text = "..." replaces all content.
        cell.text = header_text
        # cell.text_frame grants access to TextFrame for advanced formatting (paragraphs, runs).
        text_frame = cell.text_frame
        # Enable word wrapping to wrap long text within the cell bounds.
        text_frame.word_wrap = True
        # Iterate through paragraphs and runs (smallest text unit) for character-level formatting.
        # A run is a contiguous sequence of text with the same formatting properties.
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                # run.font provides access to font properties: bold, size, name, color, italic, etc.
                run.font.bold = True
                # Pt(n) converts font size from points to EMUs. 9pt is small for table headers.
                run.font.size = Pt(9)

    # Data rows - populate each row with cell data. enumerate(rows, start=1) skips row 0 (header).
    for row_idx, row_data in enumerate(rows, start=1):
        # Collect the 6 cell contents from the row_data object (each is a list of SectionContent).
        cell_sections = [
            row_data.publication_and_testing,
            row_data.proband,
            row_data.variant_info,
            row_data.clinical_presentation,
            row_data.functional_segregation,
            row_data.score,
        ]
        # Populate each cell in the data row.
        for col_idx, sections in enumerate(cell_sections):
            # Access the cell at (row_idx, col_idx). Note: row 0 is always the header.
            cell = table.cell(row_idx, col_idx)
            # Obtain the TextFrame to enable advanced formatting (bold titles, content grouping).
            text_frame = cell.text_frame
            text_frame.word_wrap = True

            # Use helper function to format cell content with bold section titles and line breaks.
            _format_cell_with_sections(text_frame, sections)

            # Apply consistent 8pt font size to all text in the cell.
            # This loop iterates the generated paragraphs and runs created by the helper function.
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(8)

    # Add pedigree image if available (bottom third of first column).
    # slide.shapes.add_picture(path, left, top, width=..., height=...) inserts an image shape.
    # Position at (0.5", 5") with 2" width; height scales automatically to preserve aspect ratio.
    for row_data in rows:
        if row_data.pedigree_image_path:
            image_path = Path(row_data.pedigree_image_path)
            # Verify the image file exists before attempting to insert it.
            if image_path.exists():
                # Add the image to the slide. Path must be a string.
                # left, top, width are in EMUs (via Inches). Height auto-scales if omitted.
                slide.shapes.add_picture(
                    str(image_path), Inches(0.5), Inches(5), width=Inches(2)
                )

    # Save the presentation to a BytesIO buffer (in-memory file).
    # This allows returning the PPTX as raw bytes without writing to disk.
    pptx_bytes_io = io.BytesIO()
    # prs.save(file) writes the presentation to the given file object or path.
    # PPTX files are actually ZIP archives containing XML files; .save() handles serialization.
    prs.save(pptx_bytes_io)
    # Return the bytes from the buffer. getvalue() retrieves the complete buffer contents.
    return pptx_bytes_io.getvalue()
